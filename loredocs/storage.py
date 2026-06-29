"""
LoreDocs Storage Layer

Manages the local filesystem and SQLite database that back all vault operations.
Storage is organized as:

    ~/.loredocs/
        config.json
        loredocs.db        (SQLite with FTS5)
        vaults/
            {vault_id}/
                docs/
                    {doc_id}/
                        current{.ext}
                        metadata.json
                        extracted.txt
                        history/
                            v1{.ext}
                            v2{.ext}
"""

import hashlib
import json
import os
import re
import shutil
import sqlite3
import time
import uuid
from contextlib import contextmanager
from datetime import datetime, timezone
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple

from .tiers import TierEnforcer, TierLimitError, get_tier, TIER_PRO  # noqa: F401 (re-exported)


# ---------------------------------------------------------------------------
# Phase 2a: Deterministic link ID + per-vault circuit breaker
# ---------------------------------------------------------------------------

def _make_link_id(
    source_vault_id: str,
    source_doc_id: str,
    target_vault_id: str,
    target_doc_id: str,
    label: str,
) -> str:
    """Return a deterministic 16-char hex ID that covers the full UNIQUE key.

    The UNIQUE constraint on doc_links is (source_vault_id, source_doc_id,
    target_vault_id, target_doc_id, label). The ID is a SHA-1 hex digest of
    all five components so id collision always implies the same logical link.
    """
    raw = f"{source_vault_id}|{source_doc_id}|{target_vault_id}|{target_doc_id}|{label}"
    return hashlib.sha1(raw.encode()).hexdigest()[:16]


_EMBD_THRESHOLD = 5
_EMBD_RESET_SECONDS = 1800

# vault_id -> (failure_count, circuit_open, opened_at)
_embd_circuits: Dict[str, Tuple[int, bool, float]] = {}


def _embd_check_circuit(vault_id: str) -> bool:
    state = _embd_circuits.get(vault_id, (0, False, 0.0))
    failures, open_, opened_at = state
    if not open_:
        return True
    reset_secs = int(os.environ.get("LOREDOCS_CIRCUIT_RESET_MINUTES", "30")) * 60
    if time.time() - opened_at >= reset_secs:
        _embd_circuits[vault_id] = (0, False, 0.0)
        return True
    return False


def _embd_record_success(vault_id: str) -> None:
    _embd_circuits.pop(vault_id, None)


def _embd_record_failure(vault_id: str) -> None:
    import logging
    failures, open_, opened_at = _embd_circuits.get(vault_id, (0, False, 0.0))
    if open_:
        _embd_circuits[vault_id] = (failures + 1, True, opened_at)
        return
    failures += 1
    if failures >= _EMBD_THRESHOLD:
        _embd_circuits[vault_id] = (failures, True, time.time())
        logging.getLogger(__name__).warning(
            "auto_link embedding circuit OPEN for vault=%r after %d failures. "
            "Will retry in 30 min.",
            vault_id, _EMBD_THRESHOLD,
        )
    else:
        _embd_circuits[vault_id] = (failures, False, 0.0)
        logging.getLogger(__name__).error(
            "auto_link_doc_embeddings failure %d/%d for vault=%r",
            failures, _EMBD_THRESHOLD, vault_id,
        )

_STOPWORDS = frozenset({
    "a", "about", "after", "again", "ago", "all", "also", "an", "and",
    "any", "are", "as", "at", "be", "been", "being", "but", "by", "can",
    "could", "did", "do", "does", "done", "down", "during", "each",
    "either", "every", "few", "for", "from", "get", "got", "had", "has",
    "have", "he", "her", "here", "him", "his", "how", "i", "if", "in",
    "into", "is", "it", "its", "just", "like", "may", "me", "more",
    "most", "my", "neither", "new", "no", "nor", "not", "now", "of",
    "off", "on", "once", "only", "or", "other", "our", "out", "over",
    "own", "per", "run", "s", "set", "she", "so", "some", "such", "t",
    "than", "that", "the", "their", "them", "then", "there", "these",
    "they", "this", "those", "through", "to", "too", "under", "up",
    "use", "used", "was", "we", "were", "what", "when", "where", "which",
    "while", "who", "will", "with", "would", "yet", "you", "your",
})


def _extract_keywords(text: str, top_n: int = 30) -> List[Tuple[str, int]]:
    """Extract top-N keywords from text using simple term frequency.

    Returns list of (term, frequency) tuples sorted descending by frequency.
    Only includes tokens of 3+ characters not in the stopword list.
    """
    if not text:
        return []
    tokens = re.findall(r'[a-z]{3,}', text.lower())
    counts: Dict[str, int] = {}
    for token in tokens:
        if token not in _STOPWORDS:
            counts[token] = counts.get(token, 0) + 1
    return sorted(counts.items(), key=lambda x: x[1], reverse=True)[:top_n]


# ---------------------------------------------------------------------------
# Constants
# ---------------------------------------------------------------------------

DEFAULT_ROOT = Path.home() / ".loredocs"
CONFIG_FILE = "config.json"
DB_FILE = "loredocs.db"
VAULTS_DIR = "vaults"

# Cross-product linking (Phase 2b, SH-10727)
# Increment when cross_product_links schema changes incompatibly.
CROSS_LINK_SCHEMA_VERSION = 1
# LoreConvo callers check >= REQUIRED_CROSS_LINK_SCHEMA_VERSION before using links.
REQUIRED_CROSS_LINK_SCHEMA_VERSION = 1
_CROSS_LINK_THRESHOLD = 0.80   # cosine similarity floor
_CROSS_LINK_L2_THRESHOLD = 0.6  # L2-normalized distance ceiling (cosine 0.80 ~ dist 0.632)
_CROSS_LINK_CAP = 5             # max cross-product links per session/doc
_CROSS_LINK_DEBOUNCE_SECS = 600 # 10-minute per-entity debounce
_CROSS_LINK_EMBEDDING_MODEL = "BAAI/bge-small-en-v1.5"
_CROSS_LINK_EMBEDDING_DIM = 384

# File size limits (bytes)
MAX_FILE_SIZE = 30 * 1024 * 1024  # 30 MB per file (matches Claude Projects)


# ---------------------------------------------------------------------------
# Cross-product path discovery (Phase 2b)
# ---------------------------------------------------------------------------

class DiscoveryError(Exception):
    """Raised when an env-var DB path override is set but invalid."""


def discover_product_db(product: str) -> Optional[Path]:
    """Return the DB path for another Lore product, or None if not installed.

    Checks the <PRODUCT>_DB_PATH env var first (for testing / non-default install).
    Validated env paths must be under $HOME, end in .db, and be readable SQLite.

    Raises DiscoveryError only when the env var is explicitly set but the path
    fails validation (misconfiguration signal worth surfacing). Returns None
    silently when the product is simply not installed.
    """
    import logging as _logging
    log = _logging.getLogger(__name__)

    env_key = f"{product.upper().replace('-', '_')}_DB_PATH"
    env_override = os.environ.get(env_key)
    if env_override:
        p = Path(env_override)
        home = Path.home()
        try:
            p.resolve().relative_to(home.resolve())
        except ValueError:
            raise DiscoveryError(
                f"{env_key}={env_override!r} is outside $HOME -- refusing to open it"
            )
        if p.suffix != ".db":
            raise DiscoveryError(
                f"{env_key}={env_override!r} does not end in .db -- check the path"
            )
        if not p.exists():
            raise DiscoveryError(
                f"{env_key}={env_override!r} set but file does not exist. "
                f"Check {product} installation or unset {env_key}."
            )
        # Validate it is a readable SQLite file
        try:
            import sqlite3 as _sqlite3
            c = _sqlite3.connect(f"file:{p}?mode=ro", uri=True)
            c.execute("SELECT 1")
            c.close()
        except Exception as exc:
            raise DiscoveryError(
                f"{env_key}={env_override!r} is not a readable SQLite file: {exc}"
            )
        return p

    default = Path.home() / f".{product}" / f"{product}.db"
    if default.exists():
        log.debug("discover_product_db: found %s at default path", product)
        return default
    log.debug("discover_product_db: %s not installed (default path absent)", product)
    return None

# Free tier limits
FREE_MAX_VAULTS = 3
FREE_MAX_DOCS_PER_VAULT = 50
FREE_MAX_STORAGE_BYTES = 500 * 1024 * 1024  # 500 MB total
FREE_MAX_VERSIONS = 5


# ---------------------------------------------------------------------------
# Text extraction helpers
# ---------------------------------------------------------------------------

def extract_text(file_path: Path) -> str:
    """Extract searchable text from a file based on its extension.

    For text-native formats, reads the file directly.
    For rich documents (PDF, DOCX, XLSX, PPTX), uses appropriate libraries.
    For binary/image/audio files, returns empty string (no extraction in v1).
    """
    suffix = file_path.suffix.lower()

    # Text-native formats -- read directly
    if suffix in (".txt", ".md", ".csv", ".html", ".htm", ".json", ".xml",
                   ".py", ".js", ".ts", ".sql", ".yaml", ".yml", ".toml",
                   ".ini", ".cfg", ".conf", ".log", ".sh", ".bat", ".ps1"):
        try:
            return file_path.read_text(encoding="utf-8", errors="replace")
        except Exception:
            return ""

    # PDF
    if suffix == ".pdf":
        try:
            import pdfplumber
            text_parts = []
            with pdfplumber.open(file_path) as pdf:
                for page in pdf.pages:
                    page_text = page.extract_text()
                    if page_text:
                        text_parts.append(page_text)
            return "\n\n".join(text_parts)
        except Exception:
            return ""

    # DOCX
    if suffix == ".docx":
        try:
            import docx
            doc = docx.Document(str(file_path))
            return "\n\n".join(p.text for p in doc.paragraphs if p.text.strip())
        except Exception:
            return ""

    # XLSX
    if suffix == ".xlsx":
        try:
            import openpyxl
            wb = openpyxl.load_workbook(str(file_path), read_only=True, data_only=True)
            text_parts = []
            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                text_parts.append(f"[Sheet: {sheet_name}]")
                for row in ws.iter_rows(values_only=True):
                    cells = [str(c) if c is not None else "" for c in row]
                    if any(cells):
                        text_parts.append("\t".join(cells))
            wb.close()
            return "\n".join(text_parts)
        except Exception:
            return ""

    # PPTX
    if suffix == ".pptx":
        try:
            from pptx import Presentation
            prs = Presentation(str(file_path))
            text_parts = []
            for i, slide in enumerate(prs.slides, 1):
                text_parts.append(f"[Slide {i}]")
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            text = paragraph.text.strip()
                            if text:
                                text_parts.append(text)
            return "\n\n".join(text_parts)
        except Exception:
            return ""

    # Images, audio, and unknown formats -- no extraction in v1
    return ""


# ---------------------------------------------------------------------------
# Database setup
# ---------------------------------------------------------------------------

def _is_in_memory_db(db_path) -> bool:
    """True for in-memory SQLite databases.

    In-memory databases cannot use WAL -- ``PRAGMA journal_mode=WAL`` always
    returns "memory" -- and there is no file for a second client to lock into a
    conflicting mode, so the WAL mixing guard does not apply to them. The guard
    keys on the configured path (not the returned mode) so a file DB that fails
    WAL still raises.
    """
    p = str(db_path)
    if p == ":memory:":
        return True
    if p.startswith("file:"):
        return ":memory:" in p or "mode=memory" in p
    return False


def _init_db(db_path: Path) -> None:
    """Create the SQLite database with FTS5 tables if they don't exist."""
    conn = sqlite3.connect(str(db_path))
    row = conn.execute("PRAGMA journal_mode=WAL").fetchone()
    actual_mode = row[0] if row else "unknown"
    if actual_mode != "wal" and not _is_in_memory_db(db_path):
        conn.close()
        raise RuntimeError(
            f"Database at '{db_path}' is in '{actual_mode}' journal mode, expected WAL. "
            "Another process may be using a conflicting journal mode. "
            "Close all other LoreDocs connections and retry."
        )
    conn.execute("PRAGMA busy_timeout=5000")
    conn.execute("PRAGMA foreign_keys=ON")

    conn.executescript("""
        CREATE TABLE IF NOT EXISTS vaults (
            id TEXT PRIMARY KEY,
            name TEXT NOT NULL,
            description TEXT DEFAULT '',
            created_at TEXT NOT NULL,
            updated_at TEXT NOT NULL,
            archived INTEGER DEFAULT 0,
            tags TEXT DEFAULT '[]',
            linked_projects TEXT DEFAULT '[]'
        );

        CREATE TABLE IF NOT EXISTS documents (
            id TEXT PRIMARY KEY,
            vault_id TEXT NOT NULL,
            name TEXT NOT NULL,
            original_filename TEXT NOT NULL,
            file_extension TEXT DEFAULT '',
            category TEXT DEFAULT 'general',
            priority TEXT DEFAULT 'normal',
            tags TEXT DEFAULT '[]',
            notes TEXT DEFAULT '',
            created_at TEXT NOT NULL,
            updated_at TEXT NOT NULL,
            file_size_bytes INTEGER DEFAULT 0,
            version_count INTEGER DEFAULT 1,
            deleted INTEGER DEFAULT 0,
            FOREIGN KEY (vault_id) REFERENCES vaults(id) ON DELETE CASCADE
        );

        CREATE TABLE IF NOT EXISTS doc_links (
            id TEXT PRIMARY KEY,
            source_vault_id TEXT NOT NULL,
            source_doc_id TEXT NOT NULL,
            target_vault_id TEXT NOT NULL,
            target_doc_id TEXT NOT NULL,
            created_at TEXT NOT NULL,
            FOREIGN KEY (source_vault_id) REFERENCES vaults(id),
            FOREIGN KEY (target_vault_id) REFERENCES vaults(id)
        );

        CREATE VIRTUAL TABLE IF NOT EXISTS doc_fts USING fts5(
            doc_id,
            vault_id,
            name,
            content,
            tags,
            notes,
            content_rowid='rowid'
        );

        CREATE INDEX IF NOT EXISTS idx_documents_vault ON documents(vault_id);
        CREATE INDEX IF NOT EXISTS idx_documents_deleted ON documents(deleted);
        CREATE INDEX IF NOT EXISTS idx_documents_category ON documents(category);
        CREATE INDEX IF NOT EXISTS idx_doc_links_target ON doc_links(target_doc_id);

        -- FTS sync triggers: keep doc_fts in sync with documents for UPDATE and DELETE.
        -- INSERT sync is handled by application code (content column requires extracted
        -- text from disk that is not available in the documents table row).
        CREATE TRIGGER IF NOT EXISTS documents_au AFTER UPDATE ON documents BEGIN
            UPDATE doc_fts SET
                vault_id = new.vault_id,
                name = new.name,
                tags = new.tags,
                notes = new.notes
            WHERE doc_id = new.id;
        END;

        CREATE TRIGGER IF NOT EXISTS documents_ad AFTER DELETE ON documents BEGIN
            DELETE FROM doc_fts WHERE doc_id = old.id;
        END;
    """)

    conn.commit()
    conn.close()


# ---------------------------------------------------------------------------
# Schema migrations (safe to run on every startup)
# ---------------------------------------------------------------------------

def _migrate_db(db_path: Path) -> None:
    """Apply incremental schema migrations that are safe to run repeatedly."""
    conn = sqlite3.connect(str(db_path))
    conn.row_factory = sqlite3.Row
    conn.execute("PRAGMA busy_timeout=5000")
    conn.execute("PRAGMA foreign_keys=ON")
    row = conn.execute("PRAGMA journal_mode").fetchone()
    actual_mode = row[0] if row else "unknown"
    if actual_mode != "wal" and not _is_in_memory_db(db_path):
        conn.close()
        raise RuntimeError(
            f"Database at '{db_path}' is in '{actual_mode}' journal mode, expected WAL."
        )

    # v0.2: add label column to doc_links if missing
    existing_cols = {row[1] for row in conn.execute("PRAGMA table_info(doc_links)")}
    if "label" not in existing_cols:
        conn.execute("ALTER TABLE doc_links ADD COLUMN label TEXT DEFAULT 'related'")

    # v0.3: add FTS sync triggers for doc_fts (MEG-00073)
    conn.executescript("""
        CREATE TRIGGER IF NOT EXISTS documents_au AFTER UPDATE ON documents BEGIN
            UPDATE doc_fts SET
                vault_id = new.vault_id,
                name = new.name,
                tags = new.tags,
                notes = new.notes
            WHERE doc_id = new.id;
        END;

        CREATE TRIGGER IF NOT EXISTS documents_ad AFTER DELETE ON documents BEGIN
            DELETE FROM doc_fts WHERE doc_id = old.id;
        END;
    """)

    # v0.4: add doc_cooccurrences table for relationship layer (RON-00067)
    conn.execute("""
        CREATE TABLE IF NOT EXISTS doc_cooccurrences (
            term    TEXT NOT NULL,
            doc_id  TEXT NOT NULL,
            frequency INTEGER NOT NULL DEFAULT 1,
            PRIMARY KEY (term, doc_id)
        )
    """)
    conn.execute(
        "CREATE INDEX IF NOT EXISTS idx_doc_cooccurrences_doc ON doc_cooccurrences(doc_id)"
    )
    conn.execute(
        "CREATE INDEX IF NOT EXISTS idx_doc_cooccurrences_term ON doc_cooccurrences(term)"
    )

    # Populate doc_cooccurrences for existing documents if table is empty
    count = conn.execute("SELECT COUNT(*) FROM doc_cooccurrences").fetchone()[0]
    if count == 0:
        _reindex_all_docs(conn, db_path.parent / VAULTS_DIR)

    # v0.5: add workspace_path to vaults for directory-scoped auto-vault (GINA-00018)
    vault_cols = {row[1] for row in conn.execute("PRAGMA table_info(vaults)")}
    if "workspace_path" not in vault_cols:
        conn.execute("ALTER TABLE vaults ADD COLUMN workspace_path TEXT DEFAULT NULL")
    conn.execute(
        "CREATE INDEX IF NOT EXISTS idx_vaults_workspace ON vaults(workspace_path)"
    )

    # v0.6: Phase 2a embedding autolink schema (SH-10529)
    # Add archived_at column to doc_links (row-level visibility for downgraded users)
    link_cols = {row[1] for row in conn.execute("PRAGMA table_info(doc_links)")}
    if "archived_at" not in link_cols:
        conn.execute("ALTER TABLE doc_links ADD COLUMN archived_at TEXT DEFAULT NULL")
    # Dedup any pre-existing auto-generated link duplicates before creating the unique index.
    # Manual links are never touched. Order matters: dedupe MUST run before CREATE UNIQUE INDEX,
    # otherwise SQLite raises IntegrityError on the existing duplicate rows.
    conn.execute("""
        DELETE FROM doc_links
        WHERE label LIKE 'auto:%'
          AND id NOT IN (
            SELECT MIN(id) FROM doc_links
            WHERE label LIKE 'auto:%'
            GROUP BY source_vault_id, source_doc_id, target_vault_id, target_doc_id, label
          )
    """)
    # SQLite doesn't support ADD CONSTRAINT on existing tables; recreate is too disruptive.
    # Use a unique index instead -- enforces the same constraint without recreation.
    conn.execute("""
        CREATE UNIQUE INDEX IF NOT EXISTS idx_doc_links_unique_key
        ON doc_links(source_vault_id, source_doc_id, target_vault_id, target_doc_id, label)
    """)

    # v0.7: Phase 2b cross-product linking (SH-10727)
    conn.executescript("""
        CREATE TABLE IF NOT EXISTS cross_product_links (
            id              INTEGER PRIMARY KEY AUTOINCREMENT,
            source_product  TEXT NOT NULL,
            source_id       TEXT NOT NULL,
            target_product  TEXT NOT NULL,
            target_id       TEXT NOT NULL,
            similarity_score REAL,
            embedding_model  TEXT NOT NULL,
            embedding_dim    INTEGER NOT NULL,
            link_type        TEXT NOT NULL DEFAULT 'auto',
            tier_required    TEXT NOT NULL DEFAULT 'pro',
            created_at       TEXT NOT NULL DEFAULT (strftime('%Y-%m-%dT%H:%M:%SZ', 'now')),
            UNIQUE(source_product, source_id, target_product, target_id)
        );
        CREATE INDEX IF NOT EXISTS idx_cross_links_source
            ON cross_product_links(source_product, source_id);
        CREATE INDEX IF NOT EXISTS idx_cross_links_target
            ON cross_product_links(target_product, target_id);
    """)

    # Privacy: vault-level opt-out column
    vault_cols_v7 = {row[1] for row in conn.execute("PRAGMA table_info(vaults)")}
    if "cross_link_opt_out" not in vault_cols_v7:
        conn.execute(
            "ALTER TABLE vaults ADD COLUMN cross_link_opt_out INTEGER NOT NULL DEFAULT 0"
        )

    # Debounce: last_cross_linked_at on documents
    doc_cols_v7 = {row[1] for row in conn.execute("PRAGMA table_info(documents)")}
    if "last_cross_linked_at" not in doc_cols_v7:
        conn.execute(
            "ALTER TABLE documents ADD COLUMN last_cross_linked_at TEXT DEFAULT NULL"
        )

    conn.commit()
    conn.close()


def _insert_doc_cooccurrences(
    conn: sqlite3.Connection,
    doc_id: str,
    name: str,
    tags: List[str],
    content: str,
) -> None:
    """(Re)index keyword co-occurrences for one document. Replaces prior index."""
    conn.execute("DELETE FROM doc_cooccurrences WHERE doc_id = ?", (doc_id,))
    text = " ".join(filter(None, [name, content]))
    for term, freq in _extract_keywords(text):
        conn.execute(
            "INSERT OR REPLACE INTO doc_cooccurrences (term, doc_id, frequency) VALUES (?, ?, ?)",
            (term, doc_id, freq)
        )
    # Boost tags as explicit terms
    for tag in tags:
        if isinstance(tag, str):
            tag_term = re.sub(r'[^a-z]', '', tag.lower())
            if len(tag_term) >= 3 and tag_term not in _STOPWORDS:
                conn.execute(
                    "INSERT OR REPLACE INTO doc_cooccurrences (term, doc_id, frequency) VALUES (?, ?, ?)",
                    (tag_term, doc_id, 5)
                )


def _auto_link_doc_cooccurrences(
    conn: sqlite3.Connection,
    doc_id: str,
    min_shared_terms: int = 3,
) -> None:
    """Create bidirectional auto-links in doc_links for docs sharing >= min_shared_terms terms.

    Uses label='auto:cooccurrence'. INSERT OR IGNORE preserves manual links.
    Caps at 20 new links per document.
    """
    src_row = conn.execute(
        "SELECT vault_id FROM documents WHERE id = ?", (doc_id,)
    ).fetchone()
    if not src_row:
        return
    source_vault_id = src_row[0]

    candidates = conn.execute(
        """SELECT dc2.doc_id, COUNT(*) as shared_count
           FROM doc_cooccurrences dc1
           JOIN doc_cooccurrences dc2 ON dc1.term = dc2.term
           WHERE dc1.doc_id = ?
             AND dc2.doc_id != ?
           GROUP BY dc2.doc_id
           HAVING COUNT(*) >= ?
           ORDER BY shared_count DESC
           LIMIT 20""",
        (doc_id, doc_id, min_shared_terms)
    ).fetchall()

    now = datetime.now(timezone.utc).isoformat()
    for row in candidates:
        other_doc_id = row[0]
        tgt_row = conn.execute(
            "SELECT vault_id FROM documents WHERE id = ? AND deleted = 0",
            (other_doc_id,)
        ).fetchone()
        if not tgt_row:
            continue
        target_vault_id = tgt_row[0]
        try:
            conn.execute(
                """INSERT OR IGNORE INTO doc_links
                   (id, source_vault_id, source_doc_id, target_vault_id, target_doc_id,
                    created_at, label)
                   VALUES (?, ?, ?, ?, ?, ?, 'auto:cooccurrence')""",
                (str(uuid.uuid4())[:12], source_vault_id, doc_id,
                 target_vault_id, other_doc_id, now)
            )
            conn.execute(
                """INSERT OR IGNORE INTO doc_links
                   (id, source_vault_id, source_doc_id, target_vault_id, target_doc_id,
                    created_at, label)
                   VALUES (?, ?, ?, ?, ?, ?, 'auto:cooccurrence')""",
                (str(uuid.uuid4())[:12], target_vault_id, other_doc_id,
                 source_vault_id, doc_id, now)
            )
        except sqlite3.IntegrityError:
            pass


def _auto_link_doc_embeddings(
    conn: sqlite3.Connection,
    doc_id: str,
    lance_index: Any,
    cap: int = 10,
    is_pro: bool = True,
) -> None:
    """Create bidirectional embedding-based links in doc_links. Pro tier only.

    Queries the Lance chunk index for semantically similar documents in the same
    vault. Confirms vault membership against SQLite before writing (Lance metadata
    is advisory). Uses INSERT WHERE EXISTS to atomically validate both doc IDs at
    write time. Caps at cap bidirectional pairs per document.
    """
    if not is_pro:
        return
    if os.environ.get("LOREDOCS_EMBEDDING_LINKS", "1") == "0":
        return

    src_row = conn.execute(
        "SELECT vault_id FROM documents WHERE id = ? AND deleted = 0", (doc_id,)
    ).fetchone()
    if not src_row:
        return
    source_vault_id = src_row[0]

    if not _embd_check_circuit(source_vault_id):
        return

    try:
        table = lance_index._open_table()
        model = lance_index._get_model()
        # Use first 256 tokens of extracted text for the query (same as semantic search)
        extracted_row = conn.execute(
            "SELECT name FROM documents WHERE id = ?", (doc_id,)
        ).fetchone()
        query_text = extracted_row[0] if extracted_row else ""
        q_vec = model.encode([query_text])[0].tolist()

        raw = table.search(
            q_vec,
            vector_column_name="vector",
            query_type="vector",
        ).limit(50).to_list()

        # Dedup to best distance per doc_id, filter by threshold + same vault
        best: Dict[str, float] = {}
        for r in raw:
            cand_doc_id = r.get("doc_id") or r.get("session_id")
            cand_vault_id = r.get("vault_id")
            dist = r.get("_distance", 999)
            if cand_doc_id == doc_id:
                continue
            if dist > 0.707:  # cosine < 0.75 for L2-normalized vectors
                continue
            if cand_vault_id != source_vault_id:
                continue
            if cand_doc_id not in best or dist < best[cand_doc_id]:
                best[cand_doc_id] = dist

        # SQLite confirmation: verify vault membership against DB (Lance metadata advisory)
        confirmed: List[str] = []
        for cand_doc_id in best:
            row = conn.execute(
                "SELECT vault_id FROM documents WHERE id = ? AND deleted = 0",
                (cand_doc_id,)
            ).fetchone()
            if row and row[0] == source_vault_id:
                confirmed.append(cand_doc_id)

        now = datetime.now(timezone.utc).isoformat()
        inserted = 0
        for other_doc_id in confirmed:
            if inserted >= cap:
                break
            existing = conn.execute(
                "SELECT 1 FROM doc_links WHERE "
                "(source_doc_id=? AND target_doc_id=?) OR "
                "(source_doc_id=? AND target_doc_id=?)",
                (doc_id, other_doc_id, other_doc_id, doc_id)
            ).fetchone()
            if existing:
                continue

            fwd_id = _make_link_id(source_vault_id, doc_id, source_vault_id, other_doc_id, "auto:embedding")
            rev_id = _make_link_id(source_vault_id, other_doc_id, source_vault_id, doc_id, "auto:embedding")
            try:
                # INSERT WHERE EXISTS validates both docs exist atomically
                conn.execute(
                    """INSERT OR IGNORE INTO doc_links
                       (id, source_vault_id, source_doc_id, target_vault_id,
                        target_doc_id, created_at, label)
                       SELECT ?, ?, ?, ?, ?, ?, 'auto:embedding'
                       WHERE EXISTS (SELECT 1 FROM documents WHERE id=? AND vault_id=? AND deleted=0)
                         AND EXISTS (SELECT 1 FROM documents WHERE id=? AND vault_id=? AND deleted=0)
                    """,
                    (fwd_id, source_vault_id, doc_id, source_vault_id, other_doc_id, now,
                     doc_id, source_vault_id, other_doc_id, source_vault_id)
                )
                conn.execute(
                    """INSERT OR IGNORE INTO doc_links
                       (id, source_vault_id, source_doc_id, target_vault_id,
                        target_doc_id, created_at, label)
                       SELECT ?, ?, ?, ?, ?, ?, 'auto:embedding'
                       WHERE EXISTS (SELECT 1 FROM documents WHERE id=? AND vault_id=? AND deleted=0)
                         AND EXISTS (SELECT 1 FROM documents WHERE id=? AND vault_id=? AND deleted=0)
                    """,
                    (rev_id, source_vault_id, other_doc_id, source_vault_id, doc_id, now,
                     other_doc_id, source_vault_id, doc_id, source_vault_id)
                )
                inserted += 1
            except sqlite3.IntegrityError:
                pass
        _embd_record_success(source_vault_id)
    except Exception as exc:
        _embd_record_failure(source_vault_id)
        import logging
        logging.getLogger(__name__).error(
            "auto_link_doc_embeddings failed for doc %s: %s", doc_id, exc
        )


def _reindex_all_docs(conn: sqlite3.Connection, vaults_dir: Path) -> None:
    """Populate doc_cooccurrences for all existing documents (first-run migration)."""
    rows = conn.execute(
        "SELECT id, vault_id, name, tags FROM documents WHERE deleted = 0"
    ).fetchall()
    for row in rows:
        doc_id = row["id"]
        vault_id = row["vault_id"]
        name = row["name"]
        tags = json.loads(row["tags"] or "[]")

        extracted_path = vaults_dir / vault_id / "docs" / doc_id / "extracted.txt"
        content = ""
        if extracted_path.exists():
            try:
                content = extracted_path.read_text(encoding="utf-8", errors="replace")
            except Exception:
                pass

        _insert_doc_cooccurrences(conn, doc_id, name, tags, content)


def _extract_frontmatter_tags(content: bytes) -> List[str]:
    """Extract tags from YAML frontmatter in a markdown file (stdlib only, no PyYAML)."""
    try:
        text = content.decode("utf-8", errors="ignore")
        if not text.startswith("---"):
            return []
        end = text.find("\n---", 3)
        if end == -1:
            return []
        block = text[3:end]
        in_tags = False
        tags: List[str] = []
        for line in block.splitlines():
            if line.startswith("tags:"):
                raw = line[5:].strip()
                if raw.startswith("[") and raw.endswith("]"):
                    return [t.strip().strip("'\"") for t in raw[1:-1].split(",") if t.strip()]
                if raw:
                    return [raw.strip()]
                in_tags = True
            elif in_tags:
                stripped = line.strip()
                if stripped.startswith("- "):
                    tags.append(stripped[2:].strip().strip("'\""))
                elif stripped and not stripped.startswith("#"):
                    break
        return tags
    except Exception:
        return []


# ---------------------------------------------------------------------------
# Storage class
# ---------------------------------------------------------------------------

def _parse_json_list(value) -> list:
    """Parse a JSON list field that may be None, empty string, or valid JSON."""
    if not value:
        return []
    try:
        result = json.loads(value)
        return result if isinstance(result, list) else []
    except (json.JSONDecodeError, TypeError):
        return []


class VaultStorage:
    """Manages the local filesystem and SQLite database for LoreDocs."""

    def __init__(self, root: Optional[Path] = None):
        self.root = root or DEFAULT_ROOT
        self.db_path = self.root / DB_FILE
        self.vaults_dir = self.root / VAULTS_DIR

        # Ensure directories exist
        self.root.mkdir(parents=True, exist_ok=True)
        self.vaults_dir.mkdir(parents=True, exist_ok=True)

        # Initialize database
        _init_db(self.db_path)
        _migrate_db(self.db_path)

        # Tier enforcer (reads config.json from root)
        self.enforcer = TierEnforcer(self.root)

        # Lance index: lazy init, Pro only
        self._lance_index = None

    @contextmanager
    def _db(self):
        """Context manager for database connections."""
        conn = sqlite3.connect(str(self.db_path))
        conn.row_factory = sqlite3.Row
        # Wait up to 5s for a competing writer instead of failing instantly
        # with "database is locked" under transient contention.
        conn.execute("PRAGMA busy_timeout=5000")
        conn.execute("PRAGMA foreign_keys=ON")
        try:
            yield conn
            conn.commit()
        except Exception:
            conn.rollback()
            raise
        finally:
            conn.close()

    def _now(self) -> str:
        """Current UTC timestamp as ISO string."""
        return datetime.now(timezone.utc).isoformat()

    # -------------------------------------------------------------------
    # LanceDB hybrid search (Pro tier)
    # -------------------------------------------------------------------

    def _get_lance_index(self):
        """Return the DocLanceIndex instance, creating it lazily on first call."""
        if self._lance_index is None:
            from .semantic_search import DocLanceIndex
            lance_dir = self.root / 'docs.lance'
            self._lance_index = DocLanceIndex(lance_dir)
        return self._lance_index

    def _lance_write_safe(self, doc_id: str, vault_id: str, name: str, text: str) -> None:
        """Index a document in Lance (Pro only). Errors are logged, never raised."""
        if get_tier(self.root) != TIER_PRO:
            return
        try:
            self._get_lance_index().index_document(doc_id, vault_id, name, text)
        except Exception as exc:
            import logging
            logging.getLogger(__name__).error(
                "Lance write failed for doc %s: %s", doc_id, exc
            )

    def _lance_delete_safe(self, doc_id: str) -> None:
        """Remove a document's chunks from Lance (best-effort)."""
        if get_tier(self.root) != TIER_PRO:
            return
        try:
            self._get_lance_index().delete_document(doc_id)
        except Exception as exc:
            import logging
            logging.getLogger(__name__).error(
                "Lance delete failed for doc %s: %s", doc_id, exc
            )

    def search_semantic(
        self,
        query: str,
        vault_id: Optional[str] = None,
        limit: int = 20,
    ) -> Dict[str, Any]:
        """Hybrid semantic search via LanceDB (Pro only).

        Returns the same dict structure as search(). Falls back to FTS5 if
        the Lance index is unavailable (index not built or search error).
        Callers must check tier before calling; this does not enforce tier.
        """
        doc_ids = self._get_lance_index().search(query, vault_id=vault_id, limit=limit)
        if not doc_ids:
            return self.search(query, vault_id=vault_id, limit=limit)

        results = []
        with self._db() as conn:
            for doc_id in doc_ids:
                row = conn.execute(
                    """SELECT d.id, d.vault_id, d.name, d.file_size_bytes,
                              v.name as vault_name
                       FROM documents d
                       JOIN vaults v ON d.vault_id = v.id
                       WHERE d.id = ? AND d.deleted = 0""",
                    (doc_id,)
                ).fetchone()
                if not row:
                    continue
                if vault_id and row["vault_id"] != vault_id:
                    continue

                # Read a snippet from extracted text
                doc_dir = self.vaults_dir / row["vault_id"] / "docs" / doc_id
                extracted_path = doc_dir / "extracted.txt"
                snippet = ""
                if extracted_path.exists():
                    try:
                        snippet = extracted_path.read_text(encoding="utf-8")[:200].strip()
                    except Exception:
                        pass

                results.append({
                    "doc_id": doc_id,
                    "vault_id": row["vault_id"],
                    "vault_name": row["vault_name"],
                    "doc_name": row["name"],
                    "snippet": snippet,
                    "relevance_rank": None,
                })

        return {
            "query": query,
            "scope": vault_id or "all_vaults",
            "count": len(results),
            "results": results,
            "semantic": True,
        }

    def rebuild_lance_index(self) -> Dict[str, Any]:
        """Rebuild the LanceDB index from all non-deleted documents. Pro only.

        Reads extracted.txt for each document. Returns a summary dict.
        """
        docs = []
        with self._db() as conn:
            rows = conn.execute(
                "SELECT id, vault_id, name FROM documents WHERE deleted = 0"
            ).fetchall()
            for row in rows:
                doc_dir = self.vaults_dir / row["vault_id"] / "docs" / row["id"]
                extracted_path = doc_dir / "extracted.txt"
                if not extracted_path.exists():
                    continue
                try:
                    text = extracted_path.read_text(encoding="utf-8")
                except Exception:
                    continue
                if text.strip():
                    docs.append({
                        "doc_id": row["id"],
                        "vault_id": row["vault_id"],
                        "name": row["name"],
                        "text": text,
                    })

        chunk_count = self._get_lance_index().rebuild(docs)
        return {
            "docs_indexed": len(docs),
            "chunks_indexed": chunk_count,
        }

    # -------------------------------------------------------------------
    # Vault operations
    # -------------------------------------------------------------------

    def get_total_storage_bytes(self) -> int:
        """Return total bytes used across all non-deleted documents."""
        with self._db() as conn:
            row = conn.execute(
                "SELECT COALESCE(SUM(file_size_bytes), 0) as total FROM documents WHERE deleted = 0"
            ).fetchone()
            return int(row["total"])

    def create_vault(self, name: str, description: str = "",
                     tags: Optional[List[str]] = None,
                     linked_projects: Optional[List[str]] = None,
                     workspace_path: Optional[str] = None) -> Dict[str, Any]:
        """Create a new vault. Returns vault metadata dict.

        Raises TierLimitError if the Free tier vault limit would be exceeded.
        """
        # Tier check: count active (non-archived) vaults
        with self._db() as conn:
            row = conn.execute(
                "SELECT COUNT(*) as cnt FROM vaults WHERE archived = 0"
            ).fetchone()
            current_count = int(row["cnt"])
        self.enforcer.check_vault_count(current_count)

        vault_id = str(uuid.uuid4())[:12]
        now = self._now()
        tags = tags or []
        linked_projects = linked_projects or []

        vault_dir = self.vaults_dir / vault_id / "docs"
        vault_dir.mkdir(parents=True, exist_ok=True)

        with self._db() as conn:
            conn.execute(
                """INSERT INTO vaults
                   (id, name, description, created_at, updated_at, tags, linked_projects, workspace_path)
                   VALUES (?, ?, ?, ?, ?, ?, ?, ?)""",
                (vault_id, name, description, now, now,
                 json.dumps(tags), json.dumps(linked_projects), workspace_path)
            )

        return {
            "id": vault_id,
            "name": name,
            "description": description,
            "created_at": now,
            "updated_at": now,
            "tags": tags,
            "linked_projects": linked_projects,
            "workspace_path": workspace_path,
            "doc_count": 0,
            "total_size_bytes": 0,
        }

    def list_vaults(self, include_archived: bool = False) -> List[Dict[str, Any]]:
        """List all vaults with summary statistics."""
        with self._db() as conn:
            if include_archived:
                rows = conn.execute("SELECT * FROM vaults ORDER BY updated_at DESC").fetchall()
            else:
                rows = conn.execute(
                    "SELECT * FROM vaults WHERE archived = 0 ORDER BY updated_at DESC"
                ).fetchall()

            result = []
            for row in rows:
                vault_id = row["id"]
                stats = conn.execute(
                    """SELECT COUNT(*) as doc_count, COALESCE(SUM(file_size_bytes), 0) as total_size
                       FROM documents WHERE vault_id = ? AND deleted = 0""",
                    (vault_id,)
                ).fetchone()

                result.append({
                    "id": vault_id,
                    "name": row["name"],
                    "description": row["description"],
                    "created_at": row["created_at"],
                    "updated_at": row["updated_at"],
                    "archived": bool(row["archived"]),
                    "tags": _parse_json_list(row["tags"]),
                    "linked_projects": _parse_json_list(row["linked_projects"]),
                    "workspace_path": row["workspace_path"] if "workspace_path" in row.keys() else None,
                    "doc_count": stats["doc_count"],
                    "total_size_bytes": stats["total_size"],
                })
            return result

    def get_vault(self, vault_id: str) -> Optional[Dict[str, Any]]:
        """Get a single vault by ID. Returns None if not found."""
        with self._db() as conn:
            row = conn.execute("SELECT * FROM vaults WHERE id = ?", (vault_id,)).fetchone()
            if not row:
                return None

            stats = conn.execute(
                """SELECT COUNT(*) as doc_count, COALESCE(SUM(file_size_bytes), 0) as total_size
                   FROM documents WHERE vault_id = ? AND deleted = 0""",
                (vault_id,)
            ).fetchone()

            docs = conn.execute(
                """SELECT id, name, category, tags, updated_at, file_size_bytes
                   FROM documents WHERE vault_id = ? AND deleted = 0
                   ORDER BY updated_at DESC""",
                (vault_id,)
            ).fetchall()

            return {
                "id": row["id"],
                "name": row["name"],
                "description": row["description"],
                "created_at": row["created_at"],
                "updated_at": row["updated_at"],
                "archived": bool(row["archived"]),
                "tags": _parse_json_list(row["tags"]),
                "linked_projects": _parse_json_list(row["linked_projects"]),
                "workspace_path": row["workspace_path"] if "workspace_path" in row.keys() else None,
                "doc_count": stats["doc_count"],
                "total_size_bytes": stats["total_size"],
                "documents": [
                    {
                        "id": d["id"],
                        "name": d["name"],
                        "category": d["category"],
                        "tags": json.loads(d["tags"]),
                        "updated_at": d["updated_at"],
                        "file_size_bytes": d["file_size_bytes"],
                    }
                    for d in docs
                ],
            }

    def find_vault_by_name(self, name: str) -> Optional[Dict[str, Any]]:
        """Find a vault by name (case-insensitive). Returns first match."""
        with self._db() as conn:
            row = conn.execute(
                "SELECT id FROM vaults WHERE LOWER(name) = LOWER(?) AND archived = 0",
                (name,)
            ).fetchone()
            if row:
                return self.get_vault(row["id"])
            return None

    def get_vault_by_workspace_path(self, workspace_path: str) -> Optional[Dict[str, Any]]:
        """Find a non-archived vault by its exact workspace_path. Returns None if not found."""
        with self._db() as conn:
            row = conn.execute(
                "SELECT id FROM vaults WHERE workspace_path = ? AND archived = 0",
                (workspace_path,)
            ).fetchone()
            if row:
                return self.get_vault(row["id"])
            return None

    def archive_vault(self, vault_id: str) -> bool:
        """Archive a vault (soft delete). Returns True if found."""
        with self._db() as conn:
            cursor = conn.execute(
                "UPDATE vaults SET archived = 1, updated_at = ? WHERE id = ?",
                (self._now(), vault_id)
            )
            return cursor.rowcount > 0

    def delete_vault(self, vault_id: str) -> bool:
        """Permanently delete a vault and all its documents. Returns True if found."""
        vault_dir = self.vaults_dir / vault_id
        with self._db() as conn:
            # Remove FTS entries for all docs in this vault
            conn.execute("DELETE FROM doc_fts WHERE vault_id = ?", (vault_id,))
            # Remove doc_links
            conn.execute(
                "DELETE FROM doc_links WHERE source_vault_id = ? OR target_vault_id = ?",
                (vault_id, vault_id)
            )
            # Remove documents
            conn.execute("DELETE FROM documents WHERE vault_id = ?", (vault_id,))
            # Remove vault
            cursor = conn.execute("DELETE FROM vaults WHERE id = ?", (vault_id,))
            found = cursor.rowcount > 0

        if vault_dir.exists():
            shutil.rmtree(vault_dir)

        return found

    def link_project(self, vault_id: str, project_name: str) -> bool:
        """Associate a Claude Project name with a vault."""
        with self._db() as conn:
            row = conn.execute(
                "SELECT linked_projects FROM vaults WHERE id = ?", (vault_id,)
            ).fetchone()
            if not row:
                return False
            projects = _parse_json_list(row["linked_projects"])
            if project_name not in projects:
                projects.append(project_name)
                conn.execute(
                    "UPDATE vaults SET linked_projects = ?, updated_at = ? WHERE id = ?",
                    (json.dumps(projects), self._now(), vault_id)
                )
            return True

    # -------------------------------------------------------------------
    # Document operations
    # -------------------------------------------------------------------

    def add_document(self, vault_id: str, name: str, content: bytes,
                     filename: str, tags: Optional[List[str]] = None,
                     category: str = "general", priority: str = "normal",
                     notes: str = "") -> Optional[Dict[str, Any]]:
        """Add a document to a vault. Content is raw bytes.
        Returns document metadata dict, or None if vault not found.
        """
        # Verify vault exists
        vault = self.get_vault(vault_id)
        if not vault:
            return None

        # OPP-006: Validate filename -- reject path traversal, absolute paths, null bytes
        if not filename or '..' in filename or '\x00' in filename or os.path.isabs(filename):
            raise ValueError('Invalid filename: must be a plain name with no path components')
        if os.path.basename(filename) != filename:
            raise ValueError('Invalid filename: must not contain directory separators')

        # OPP-009: Enforce per-file size limit before writing anything to disk
        if len(content) > MAX_FILE_SIZE:
            raise ValueError(
                f'Document exceeds maximum file size ({MAX_FILE_SIZE // (1024 * 1024)} MB)'
            )

        # Tier checks before writing anything
        self.enforcer.check_doc_count(vault["doc_count"], vault_name=vault["name"])
        self.enforcer.check_storage(self.get_total_storage_bytes(), len(content))

        doc_id = str(uuid.uuid4())[:12]
        now = self._now()
        tags = tags or []
        ext = Path(filename).suffix.lower()

        # Write file to disk
        doc_dir = self.vaults_dir / vault_id / "docs" / doc_id
        doc_dir.mkdir(parents=True, exist_ok=True)
        (doc_dir / "history").mkdir(exist_ok=True)

        current_file = doc_dir / f"current{ext}"
        current_file.write_bytes(content)

        # Extract text for search indexing
        extracted = extract_text(current_file)
        (doc_dir / "extracted.txt").write_text(extracted, encoding="utf-8")

        # Write metadata
        meta = {
            "id": doc_id,
            "vault_id": vault_id,
            "name": name,
            "original_filename": filename,
            "file_extension": ext,
            "category": category,
            "priority": priority,
            "tags": tags,
            "notes": notes,
            "created_at": now,
            "updated_at": now,
            "file_size_bytes": len(content),
            "version_count": 1,
        }
        (doc_dir / "metadata.json").write_text(
            json.dumps(meta, indent=2), encoding="utf-8"
        )

        file_size = len(content)

        # Insert into database
        with self._db() as conn:
            conn.execute(
                """INSERT INTO documents
                   (id, vault_id, name, original_filename, file_extension,
                    category, priority, tags, notes, created_at, updated_at,
                    file_size_bytes, version_count)
                   VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?)""",
                (doc_id, vault_id, name, filename, ext, category, priority,
                 json.dumps(tags), notes, now, now, file_size, 1)
            )

            # Index in FTS
            conn.execute(
                """INSERT INTO doc_fts (doc_id, vault_id, name, content, tags, notes)
                   VALUES (?, ?, ?, ?, ?, ?)""",
                (doc_id, vault_id, name, extracted, " ".join(tags), notes)
            )

            # Update vault timestamp
            conn.execute(
                "UPDATE vaults SET updated_at = ? WHERE id = ?", (now, vault_id)
            )

            # Index co-occurrences and auto-link related docs (best-effort)
            try:
                _insert_doc_cooccurrences(conn, doc_id, name, tags, extracted)
                _auto_link_doc_cooccurrences(conn, doc_id)
            except Exception:
                pass

        self._lance_write_safe(doc_id, vault_id, name, extracted)
        # Phase 2a: embedding-based auto-link (Pro only, errors never propagate)
        is_pro = get_tier(self.root) == TIER_PRO
        if is_pro and self._lance_index is not None:
            with self._db() as conn:
                try:
                    _auto_link_doc_embeddings(conn, doc_id, self._lance_index, is_pro=True)
                except Exception:
                    pass
        # Phase 2b: save-triggered cross-product linking (Pro only, best-effort)
        try:
            self.cross_link_doc(doc_id, vault_id)
        except Exception:
            pass
        return meta

    def add_document_from_text(self, vault_id: str, name: str, text_content: str,
                               filename: Optional[str] = None,
                               tags: Optional[List[str]] = None,
                               category: str = "general", priority: str = "normal",
                               notes: str = "") -> Optional[Dict[str, Any]]:
        """Convenience method to add a text/markdown document from a string."""
        if filename is None:
            safe_name = name.replace(" ", "_").lower()
            filename = f"{safe_name}.md"
        content_bytes = text_content.encode("utf-8")
        return self.add_document(
            vault_id, name, content_bytes, filename,
            tags=tags, category=category, priority=priority, notes=notes
        )

    def update_document(self, doc_id: str, content: Optional[bytes] = None,
                        filename: Optional[str] = None,
                        name: Optional[str] = None,
                        tags: Optional[List[str]] = None,
                        category: Optional[str] = None,
                        priority: Optional[str] = None,
                        notes: Optional[str] = None) -> Optional[Dict[str, Any]]:
        """Update a document. Auto-versions the previous file if content changes.
        Returns updated metadata, or None if not found.
        """
        with self._db() as conn:
            row = conn.execute(
                "SELECT * FROM documents WHERE id = ? AND deleted = 0", (doc_id,)
            ).fetchone()
            if not row:
                return None

            vault_id = row["vault_id"]
            now = self._now()
            doc_dir = self.vaults_dir / vault_id / "docs" / doc_id
            ext = row["file_extension"]
            current_file = doc_dir / f"current{ext}"
            version_count = row["version_count"]

            # If content is changing, version the current file first
            if content is not None:
                if current_file.exists():
                    # Tier check: verify we haven't hit the version limit
                    self.enforcer.check_version_count(version_count, doc_name=row["name"])
                    history_file = doc_dir / "history" / f"v{version_count}{ext}"
                    shutil.copy2(current_file, history_file)
                    version_count += 1

                if filename:
                    new_ext = Path(filename).suffix.lower()
                    if new_ext != ext:
                        current_file.unlink(missing_ok=True)
                        ext = new_ext
                        current_file = doc_dir / f"current{ext}"

                current_file.write_bytes(content)

                # Re-extract text
                extracted = extract_text(current_file)
                (doc_dir / "extracted.txt").write_text(extracted, encoding="utf-8")

                file_size = len(content)
            else:
                file_size = row["file_size_bytes"]
                extracted = None

            # Build update query
            updates = {"updated_at": now, "version_count": version_count}
            if name is not None:
                updates["name"] = name
            if tags is not None:
                updates["tags"] = json.dumps(tags)
            if category is not None:
                updates["category"] = category
            if priority is not None:
                updates["priority"] = priority
            if notes is not None:
                updates["notes"] = notes
            if content is not None:
                updates["file_size_bytes"] = file_size
                updates["file_extension"] = ext
                if filename:
                    updates["original_filename"] = filename

            set_clause = ", ".join(f"{k} = ?" for k in updates)
            values = list(updates.values()) + [doc_id]
            conn.execute(
                f"UPDATE documents SET {set_clause} WHERE id = ?", values
            )

            # Update FTS index
            conn.execute("DELETE FROM doc_fts WHERE doc_id = ?", (doc_id,))
            final_name = name if name is not None else row["name"]
            final_tags = tags if tags is not None else _parse_json_list(row["tags"])
            final_notes = notes if notes is not None else row["notes"]
            if extracted is None:
                extracted_path = doc_dir / "extracted.txt"
                extracted = extracted_path.read_text(encoding="utf-8") if extracted_path.exists() else ""

            conn.execute(
                """INSERT INTO doc_fts (doc_id, vault_id, name, content, tags, notes)
                   VALUES (?, ?, ?, ?, ?, ?)""",
                (doc_id, vault_id, final_name, extracted,
                 " ".join(final_tags), final_notes)
            )

            # Update vault timestamp
            conn.execute(
                "UPDATE vaults SET updated_at = ? WHERE id = ?", (now, vault_id)
            )

            # Re-index co-occurrences and auto-link (best-effort)
            try:
                _insert_doc_cooccurrences(conn, doc_id, final_name, final_tags, extracted)
                _auto_link_doc_cooccurrences(conn, doc_id)
            except Exception:
                pass

        if content is not None or name is not None:
            self._lance_write_safe(doc_id, vault_id, final_name, extracted)
            # Phase 2a: embedding-based auto-link (Pro only, errors never propagate)
            is_pro = get_tier(self.root) == TIER_PRO
            if is_pro and self._lance_index is not None:
                with self._db() as conn:
                    try:
                        _auto_link_doc_embeddings(conn, doc_id, self._lance_index, is_pro=True)
                    except Exception:
                        pass
            # Phase 2b: save-triggered cross-product linking (Pro only, best-effort)
            try:
                self.cross_link_doc(doc_id, vault_id)
            except Exception:
                pass
        # Return fresh metadata
        return self.get_document(doc_id)

    def remove_document(self, doc_id: str) -> bool:
        """Soft-delete a document. Returns True if found."""
        deleted = False
        with self._db() as conn:
            cursor = conn.execute(
                "UPDATE documents SET deleted = 1, updated_at = ? WHERE id = ? AND deleted = 0",
                (self._now(), doc_id)
            )
            if cursor.rowcount > 0:
                conn.execute("DELETE FROM doc_fts WHERE doc_id = ?", (doc_id,))
                deleted = True
        if deleted:
            self._lance_delete_safe(doc_id)
        return deleted

    def get_document(self, doc_id: str) -> Optional[Dict[str, Any]]:
        """Get document metadata by ID."""
        with self._db() as conn:
            row = conn.execute(
                "SELECT * FROM documents WHERE id = ? AND deleted = 0", (doc_id,)
            ).fetchone()
            if not row:
                return None
            return {
                "id": row["id"],
                "vault_id": row["vault_id"],
                "name": row["name"],
                "original_filename": row["original_filename"],
                "file_extension": row["file_extension"],
                "category": row["category"],
                "priority": row["priority"],
                "tags": _parse_json_list(row["tags"]),
                "notes": row["notes"],
                "created_at": row["created_at"],
                "updated_at": row["updated_at"],
                "file_size_bytes": row["file_size_bytes"],
                "version_count": row["version_count"],
            }

    def get_document_content(self, doc_id: str) -> Optional[str]:
        """Read the extracted text content of a document."""
        with self._db() as conn:
            row = conn.execute(
                "SELECT vault_id FROM documents WHERE id = ? AND deleted = 0", (doc_id,)
            ).fetchone()
            if not row:
                return None
            extracted_path = self.vaults_dir / row["vault_id"] / "docs" / doc_id / "extracted.txt"
            if extracted_path.exists():
                return extracted_path.read_text(encoding="utf-8")
            return ""

    def get_document_raw_path(self, doc_id: str) -> Optional[Path]:
        """Get the path to the raw document file."""
        with self._db() as conn:
            row = conn.execute(
                "SELECT vault_id, file_extension FROM documents WHERE id = ? AND deleted = 0",
                (doc_id,)
            ).fetchone()
            if not row:
                return None
            return self.vaults_dir / row["vault_id"] / "docs" / doc_id / f"current{row['file_extension']}"

    def list_documents(self, vault_id: str, sort_by: str = "updated_at",
                       sort_order: str = "desc", category: Optional[str] = None,
                       tag: Optional[str] = None,
                       limit: int = 50, offset: int = 0) -> Dict[str, Any]:
        """List documents in a vault with filtering and sorting."""
        valid_sorts = {"name", "updated_at", "created_at", "file_size_bytes", "category"}
        if sort_by not in valid_sorts:
            sort_by = "updated_at"
        if sort_order.lower() not in ("asc", "desc"):
            sort_order = "desc"

        with self._db() as conn:
            where_parts = ["vault_id = ?", "deleted = 0"]
            params: list = [vault_id]

            if category:
                where_parts.append("category = ?")
                params.append(category)

            if tag:
                where_parts.append("tags LIKE ?")
                params.append(f'%"{tag}"%')

            where_clause = " AND ".join(where_parts)

            # Get total count
            total = conn.execute(
                f"SELECT COUNT(*) as cnt FROM documents WHERE {where_clause}", params
            ).fetchone()["cnt"]

            # Get page
            rows = conn.execute(
                f"""SELECT * FROM documents WHERE {where_clause}
                    ORDER BY {sort_by} {sort_order}
                    LIMIT ? OFFSET ?""",
                params + [limit, offset]
            ).fetchall()

            docs = []
            for row in rows:
                docs.append({
                    "id": row["id"],
                    "name": row["name"],
                    "original_filename": row["original_filename"],
                    "file_extension": row["file_extension"],
                    "category": row["category"],
                    "priority": row["priority"],
                    "tags": _parse_json_list(row["tags"]),
                    "notes": row["notes"],
                    "updated_at": row["updated_at"],
                    "file_size_bytes": row["file_size_bytes"],
                    "version_count": row["version_count"],
                })

            return {
                "total": total,
                "count": len(docs),
                "offset": offset,
                "has_more": total > offset + len(docs),
                "documents": docs,
            }

    # -------------------------------------------------------------------
    # Search
    # -------------------------------------------------------------------

    @staticmethod
    def _sanitize_fts_query(query: str) -> str:
        """OPP-010: Sanitize user input for FTS5 MATCH without changing search semantics.

        Strategy: quote each individual token so hyphens, colons, and other
        FTS5 operators inside a token are treated as literals, but multiple
        tokens are implicitly ANDed (the FTS5 default). This preserves the
        expected behavior where "data warehouse migration" matches documents
        containing all three words anywhere, not just as a consecutive phrase.
        """
        safe = query.strip()
        if not safe:
            return '""'
        tokens = safe.split()
        quoted = ['"' + t.replace('"', '') + '"' for t in tokens if t.replace('"', '')]
        return ' '.join(quoted) if quoted else '""'

    def search(self, query: str, vault_id: Optional[str] = None,
               limit: int = 20, offset: int = 0) -> Dict[str, Any]:
        """Full-text search across document contents using FTS5.

        If vault_id is provided, searches only that vault.
        Otherwise searches all vaults (cross-vault search).
        """
        # OPP-010: Sanitize FTS5 MATCH input to prevent operator injection
        safe_query = self._sanitize_fts_query(query)
        with self._db() as conn:
            if vault_id:
                rows = conn.execute(
                    """SELECT doc_id, vault_id, name,
                              snippet(doc_fts, 3, '>>>', '<<<', '...', 40) as snippet,
                              rank
                       FROM doc_fts
                       WHERE doc_fts MATCH ? AND vault_id = ?
                       ORDER BY rank
                       LIMIT ? OFFSET ?""",
                    (safe_query, vault_id, limit, offset)
                ).fetchall()
            else:
                rows = conn.execute(
                    """SELECT doc_id, vault_id, name,
                              snippet(doc_fts, 3, '>>>', '<<<', '...', 40) as snippet,
                              rank
                       FROM doc_fts
                       WHERE doc_fts MATCH ?
                       ORDER BY rank
                       LIMIT ? OFFSET ?""",
                    (safe_query, limit, offset)
                ).fetchall()

            results = []
            for row in rows:
                # Get vault name for context
                vault_row = conn.execute(
                    "SELECT name FROM vaults WHERE id = ?", (row["vault_id"],)
                ).fetchone()
                vault_name = vault_row["name"] if vault_row else "Unknown"

                results.append({
                    "doc_id": row["doc_id"],
                    "vault_id": row["vault_id"],
                    "vault_name": vault_name,
                    "doc_name": row["name"],
                    "snippet": row["snippet"],
                    "relevance_rank": row["rank"],
                })

            return {
                "query": query,
                "scope": vault_id or "all_vaults",
                "count": len(results),
                "results": results,
            }

    def search_by_tag(self, tag: str, vault_id: Optional[str] = None) -> List[Dict[str, Any]]:
        """Find all documents with a specific tag."""
        with self._db() as conn:
            if vault_id:
                rows = conn.execute(
                    """SELECT d.*, v.name as vault_name FROM documents d
                       JOIN vaults v ON d.vault_id = v.id
                       WHERE d.tags LIKE ? AND d.vault_id = ? AND d.deleted = 0
                       ORDER BY d.updated_at DESC""",
                    (f'%"{tag}"%', vault_id)
                ).fetchall()
            else:
                rows = conn.execute(
                    """SELECT d.*, v.name as vault_name FROM documents d
                       JOIN vaults v ON d.vault_id = v.id
                       WHERE d.tags LIKE ? AND d.deleted = 0
                       ORDER BY d.updated_at DESC""",
                    (f'%"{tag}"%',)
                ).fetchall()

            return [
                {
                    "id": r["id"],
                    "vault_id": r["vault_id"],
                    "vault_name": r["vault_name"],
                    "name": r["name"],
                    "category": r["category"],
                    "tags": json.loads(r["tags"]),
                    "updated_at": r["updated_at"],
                    "file_size_bytes": r["file_size_bytes"],
                }
                for r in rows
            ]

    # -------------------------------------------------------------------
    # Tagging and metadata
    # -------------------------------------------------------------------

    def tag_document(self, doc_id: str, add_tags: Optional[List[str]] = None,
                     remove_tags: Optional[List[str]] = None) -> Optional[List[str]]:
        """Add or remove tags from a document. Returns final tag list, or None if not found."""
        with self._db() as conn:
            row = conn.execute(
                "SELECT tags FROM documents WHERE id = ? AND deleted = 0", (doc_id,)
            ).fetchone()
            if not row:
                return None

            current_tags = set(_parse_json_list(row["tags"]))
            if add_tags:
                current_tags.update(add_tags)
            if remove_tags:
                current_tags -= set(remove_tags)

            final_tags = sorted(current_tags)
            conn.execute(
                "UPDATE documents SET tags = ?, updated_at = ? WHERE id = ?",
                (json.dumps(final_tags), self._now(), doc_id)
            )

            # Update FTS
            fts_row = conn.execute(
                "SELECT rowid, name, content, notes FROM doc_fts WHERE doc_id = ?", (doc_id,)
            ).fetchone()
            if fts_row:
                conn.execute("DELETE FROM doc_fts WHERE doc_id = ?", (doc_id,))
                conn.execute(
                    """INSERT INTO doc_fts (doc_id, vault_id, name, content, tags, notes)
                       VALUES (?, (SELECT vault_id FROM documents WHERE id = ?),
                               ?, ?, ?, ?)""",
                    (doc_id, doc_id, fts_row["name"], fts_row["content"],
                     " ".join(final_tags), fts_row["notes"])
                )

            return final_tags

    def bulk_tag(self, doc_ids: List[str], add_tags: Optional[List[str]] = None,
                 remove_tags: Optional[List[str]] = None) -> int:
        """Apply tag changes to multiple documents in a single transaction. Returns count of modified docs."""
        if not doc_ids:
            return 0
        with self._db() as conn:
            placeholders = ','.join('?' * len(doc_ids))
            rows = conn.execute(
                f"SELECT id, tags FROM documents WHERE id IN ({placeholders}) AND deleted = 0",
                doc_ids
            ).fetchall()
            count = 0
            now = self._now()
            for row in rows:
                current_tags = set(_parse_json_list(row["tags"]))
                if add_tags:
                    current_tags.update(add_tags)
                if remove_tags:
                    current_tags -= set(remove_tags)
                final_tags = sorted(current_tags)
                conn.execute(
                    "UPDATE documents SET tags = ?, updated_at = ? WHERE id = ?",
                    (json.dumps(final_tags), now, row["id"])
                )
                fts_row = conn.execute(
                    "SELECT name, content, notes FROM doc_fts WHERE doc_id = ?", (row["id"],)
                ).fetchone()
                if fts_row:
                    conn.execute("DELETE FROM doc_fts WHERE doc_id = ?", (row["id"],))
                    conn.execute(
                        """INSERT INTO doc_fts (doc_id, vault_id, name, content, tags, notes)
                           VALUES (?, (SELECT vault_id FROM documents WHERE id = ?),
                                   ?, ?, ?, ?)""",
                        (row["id"], row["id"], fts_row["name"], fts_row["content"],
                         " ".join(final_tags), fts_row["notes"])
                    )
                count += 1
        return count

    # -------------------------------------------------------------------
    # Version history
    # -------------------------------------------------------------------

    def get_doc_history(self, doc_id: str) -> Optional[List[Dict[str, Any]]]:
        """Get version history for a document. Returns list of versions."""
        with self._db() as conn:
            row = conn.execute(
                "SELECT vault_id, file_extension, version_count FROM documents WHERE id = ? AND deleted = 0",
                (doc_id,)
            ).fetchone()
            if not row:
                return None

            history_dir = self.vaults_dir / row["vault_id"] / "docs" / doc_id / "history"
            versions = []
            for i in range(1, row["version_count"]):
                version_file = history_dir / f"v{i}{row['file_extension']}"
                if version_file.exists():
                    stat = version_file.stat()
                    versions.append({
                        "version": i,
                        "file_size_bytes": stat.st_size,
                        "modified_at": datetime.fromtimestamp(
                            stat.st_mtime, tz=timezone.utc
                        ).isoformat(),
                    })

            # Current version
            versions.append({
                "version": row["version_count"],
                "file_size_bytes": 0,  # will be filled from db
                "modified_at": self._now(),
                "current": True,
            })

            return versions

    # -------------------------------------------------------------------
    # Cross-vault operations
    # -------------------------------------------------------------------

    def copy_document(self, doc_id: str, target_vault_id: str) -> Optional[Dict[str, Any]]:
        """Copy a document to another vault. Returns new document metadata."""
        source_doc = self.get_document(doc_id)
        if not source_doc:
            return None

        target_vault = self.get_vault(target_vault_id)
        if not target_vault:
            return None

        # Read the raw file
        raw_path = self.get_document_raw_path(doc_id)
        if not raw_path or not raw_path.exists():
            return None

        content = raw_path.read_bytes()
        return self.add_document(
            target_vault_id,
            name=source_doc["name"],
            content=content,
            filename=source_doc["original_filename"],
            tags=source_doc["tags"],
            category=source_doc["category"],
            priority=source_doc["priority"],
            notes=source_doc["notes"],
        )

    def move_document(self, doc_id: str, target_vault_id: str) -> Optional[Dict[str, Any]]:
        """Move a document to another vault. Returns new document metadata."""
        new_doc = self.copy_document(doc_id, target_vault_id)
        if new_doc:
            self.remove_document(doc_id)
        return new_doc

    # -------------------------------------------------------------------
    # Import / Export
    # -------------------------------------------------------------------

    def import_directory(self, vault_id: str, dir_path: Path,
                         tags: Optional[List[str]] = None,
                         category: str = "imported",
                         recursive: bool = True) -> List[Dict[str, Any]]:
        """Bulk import all supported files from a directory into a vault."""
        imported = []
        if not dir_path.is_dir():
            return imported

        iterator = sorted(dir_path.rglob("*")) if recursive else sorted(dir_path.iterdir())
        for file_path in iterator:
            if file_path.is_symlink():
                continue
            if file_path.is_file() and not file_path.name.startswith("."):
                if file_path.stat().st_size > MAX_FILE_SIZE:
                    continue
                try:
                    content = file_path.read_bytes()
                    doc_name = file_path.stem.replace("_", " ").replace("-", " ").title()
                    doc_tags = list(tags) if tags else []
                    if file_path.suffix.lower() == ".md":
                        fm_tags = _extract_frontmatter_tags(content)
                        for t in fm_tags:
                            if t and t not in doc_tags:
                                doc_tags.append(t)
                    result = self.add_document(
                        vault_id, name=doc_name, content=content,
                        filename=file_path.name,
                        tags=doc_tags if doc_tags else None,
                        category=category
                    )
                    if result:
                        imported.append(result)
                except Exception:
                    continue

        return imported

    def export_vault(self, vault_id: str, output_dir: Path) -> int:
        """Export all current documents from a vault to a directory.
        Returns number of files exported.
        """
        output_dir.mkdir(parents=True, exist_ok=True)
        count = 0

        with self._db() as conn:
            rows = conn.execute(
                "SELECT id, original_filename, file_extension FROM documents WHERE vault_id = ? AND deleted = 0",
                (vault_id,)
            ).fetchall()

            real_output_dir = os.path.realpath(str(output_dir))
            for row in rows:
                raw_path = self.vaults_dir / vault_id / "docs" / row["id"] / f"current{row['file_extension']}"
                if raw_path.exists():
                    # OPP-006: Sanitize original_filename to a plain basename before
                    # constructing the destination path, then confirm it stays inside
                    # output_dir (guards against stored traversal sequences).
                    safe_export_name = os.path.basename(row["original_filename"])
                    if not safe_export_name:
                        safe_export_name = f"{row['id']}{row['file_extension']}"
                    dest = output_dir / safe_export_name
                    real_dest = os.path.realpath(str(dest))
                    if not real_dest.startswith(real_output_dir + os.sep):
                        continue  # skip malicious path -- log ID only, not content
                    # Avoid filename collisions
                    counter = 1
                    while dest.exists():
                        stem = Path(safe_export_name).stem
                        dest = output_dir / f"{stem}_{counter}{row['file_extension']}"
                        counter += 1
                    shutil.copy2(raw_path, dest)
                    count += 1

        return count

    # -------------------------------------------------------------------
    # Document linking (Phase 2)
    # -------------------------------------------------------------------

    def link_doc(self, source_doc_id: str, target_doc_id: str,
                 label: str = "related") -> Optional[Dict[str, Any]]:
        """Create a link between two documents.

        Links are stored bidirectionally: one row per direction so queries
        work without UNION.  If the link already exists it is returned as-is.
        Returns the link record, or None if either document is not found.
        """
        with self._db() as conn:
            # Verify both documents exist and are not deleted
            src = conn.execute(
                "SELECT id, vault_id, name FROM documents WHERE id = ? AND deleted = 0",
                (source_doc_id,)
            ).fetchone()
            tgt = conn.execute(
                "SELECT id, vault_id, name FROM documents WHERE id = ? AND deleted = 0",
                (target_doc_id,)
            ).fetchone()
            if not src or not tgt:
                return None

            now = self._now()
            link_id = str(uuid.uuid4())
            reverse_id = str(uuid.uuid4())

            # Check if forward link already exists
            existing = conn.execute(
                "SELECT id FROM doc_links WHERE source_doc_id = ? AND target_doc_id = ?",
                (source_doc_id, target_doc_id)
            ).fetchone()
            if existing:
                return {
                    "id": existing["id"],
                    "source_doc_id": source_doc_id,
                    "source_doc_name": src["name"],
                    "target_doc_id": target_doc_id,
                    "target_doc_name": tgt["name"],
                    "label": label,
                    "already_existed": True,
                }

            conn.execute(
                """INSERT INTO doc_links (id, source_vault_id, source_doc_id,
                                          target_vault_id, target_doc_id, created_at, label)
                   VALUES (?, ?, ?, ?, ?, ?, ?)""",
                (link_id, src["vault_id"], source_doc_id,
                 tgt["vault_id"], target_doc_id, now, label)
            )
            # Store reverse direction for symmetric lookup
            conn.execute(
                """INSERT INTO doc_links (id, source_vault_id, source_doc_id,
                                          target_vault_id, target_doc_id, created_at, label)
                   VALUES (?, ?, ?, ?, ?, ?, ?)""",
                (reverse_id, tgt["vault_id"], target_doc_id,
                 src["vault_id"], source_doc_id, now, label)
            )

            return {
                "id": link_id,
                "source_doc_id": source_doc_id,
                "source_doc_name": src["name"],
                "target_doc_id": target_doc_id,
                "target_doc_name": tgt["name"],
                "label": label,
                "created_at": now,
                "already_existed": False,
            }

    def unlink_doc(self, source_doc_id: str, target_doc_id: str) -> int:
        """Remove all links between two documents (both directions).

        Returns the number of link rows deleted (0 if no link existed).
        """
        with self._db() as conn:
            cursor = conn.execute(
                """DELETE FROM doc_links
                   WHERE (source_doc_id = ? AND target_doc_id = ?)
                      OR (source_doc_id = ? AND target_doc_id = ?)""",
                (source_doc_id, target_doc_id, target_doc_id, source_doc_id)
            )
            return cursor.rowcount

    def archive_embedding_links(self) -> int:
        """Set archived_at on all auto:embedding links (Pro->Free downgrade).

        Defense-in-depth: find_related_docs already filters by archived_at IS NULL,
        but setting this flag at downgrade time makes the exclusion data-layer enforced
        and independent of per-call filter logic.
        """
        now = self._now()
        with self._db() as conn:
            cursor = conn.execute(
                "UPDATE doc_links SET archived_at = ? "
                "WHERE label = 'auto:embedding' AND archived_at IS NULL",
                (now,)
            )
            return cursor.rowcount

    def unarchive_embedding_links(self) -> int:
        """Clear archived_at on all auto:embedding links (Free->Pro re-upgrade)."""
        with self._db() as conn:
            cursor = conn.execute(
                "UPDATE doc_links SET archived_at = NULL "
                "WHERE label = 'auto:embedding' AND archived_at IS NOT NULL"
            )
            return cursor.rowcount

    def find_related_docs(self, doc_id: str) -> List[Dict[str, Any]]:
        """Return all documents linked to or from the given document.

        Each result includes the related document's metadata and the link label.
        Embedding links (label='auto:embedding') are excluded for free-tier callers
        and deduplicated: if the same pair has both co-occurrence and embedding links,
        the co-occurrence entry wins.
        """
        is_pro = get_tier(self.root) == TIER_PRO
        with self._db() as conn:
            rows = conn.execute(
                """SELECT d.id, d.vault_id, d.name, d.category, d.tags, d.notes,
                          d.updated_at, dl.label, v.name as vault_name
                   FROM doc_links dl
                   JOIN documents d ON dl.target_doc_id = d.id
                   JOIN vaults v ON d.vault_id = v.id
                   WHERE dl.source_doc_id = ?
                     AND d.deleted = 0
                     AND (dl.archived_at IS NULL)
                   ORDER BY dl.created_at DESC""",
                (doc_id,)
            ).fetchall()

            seen: Dict[str, dict] = {}
            for row in rows:
                item = dict(row)
                item["tags"] = json.loads(item.get("tags") or "[]")
                label = item.get("label", "related")
                target_id = item["id"]

                # Filter embedding links for free tier
                if label == "auto:embedding" and not is_pro:
                    continue

                # Dedup: cooccurrence wins over embedding for same pair
                if target_id in seen:
                    existing_label = seen[target_id].get("label", "")
                    if existing_label == "auto:cooccurrence" and label == "auto:embedding":
                        continue
                    if label == "auto:cooccurrence":
                        seen[target_id] = item
                else:
                    seen[target_id] = item

            return list(seen.values())

    # -------------------------------------------------------------------
    # Vault manifest export (Phase 2)
    # -------------------------------------------------------------------

    def get_vault_manifest(self, vault_id: str) -> Optional[Dict[str, Any]]:
        """Return a JSON-serializable manifest of the vault.

        Includes vault metadata, all document metadata (not file contents),
        tag index, and document count by category.  Useful for quick overviews
        and for generating marketplace showcase exports.
        """
        with self._db() as conn:
            vault = conn.execute(
                "SELECT * FROM vaults WHERE id = ? AND archived = 0",
                (vault_id,)
            ).fetchone()
            if not vault:
                return None

            vault_dict = dict(vault)
            vault_dict["tags"] = json.loads(vault_dict.get("tags") or "[]")
            vault_dict["linked_projects"] = json.loads(vault_dict.get("linked_projects") or "[]")

            docs = conn.execute(
                """SELECT id, name, original_filename, category, priority, tags,
                          notes, created_at, updated_at, file_size_bytes, version_count
                   FROM documents
                   WHERE vault_id = ? AND deleted = 0
                   ORDER BY updated_at DESC""",
                (vault_id,)
            ).fetchall()

            doc_list = []
            tag_counts: Dict[str, int] = {}
            category_counts: Dict[str, int] = {}

            for doc in docs:
                d = dict(doc)
                d["tags"] = json.loads(d.get("tags") or "[]")
                for t in d["tags"]:
                    tag_counts[t] = tag_counts.get(t, 0) + 1
                category_counts[d["category"]] = category_counts.get(d["category"], 0) + 1
                doc_list.append(d)

            # Link count (stored bidirectionally so divide by 2)
            link_row = conn.execute(
                "SELECT COUNT(*) as cnt FROM doc_links WHERE source_vault_id = ?",
                (vault_id,)
            ).fetchone()
            link_count = (link_row["cnt"] // 2) if link_row else 0

            return {
                "vault": vault_dict,
                "document_count": len(doc_list),
                "link_count": link_count,
                "category_counts": category_counts,
                "tag_counts": tag_counts,
                "documents": doc_list,
                "generated_at": self._now(),
            }

    # -------------------------------------------------------------------
    # Cross-product linking (Phase 2b, SH-10727)
    # -------------------------------------------------------------------

    def _write_cross_product_link(
        self,
        conn: sqlite3.Connection,
        source_product: str,
        source_id: str,
        target_product: str,
        target_id: str,
        similarity_score: Optional[float],
        embedding_model: str,
        embedding_dim: int,
        link_type: str = "auto",
        tier_required: str = "pro",
    ) -> None:
        """Write a single cross-product link. All cross-product writes go through here."""
        conn.execute(
            """INSERT OR IGNORE INTO cross_product_links
               (source_product, source_id, target_product, target_id,
                similarity_score, embedding_model, embedding_dim,
                link_type, tier_required)
               VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?)""",
            (source_product, source_id, target_product, target_id,
             similarity_score, embedding_model, embedding_dim,
             link_type, tier_required),
        )

    def get_cross_product_links(
        self,
        source_product: str,
        source_id: str,
        current_embedding_model: str,
        limit: int = 5,
        is_pro: bool = False,
    ) -> Dict[str, Any]:
        """Return cross-product links for a source entity.

        Filters stale-model links (marks them is_stale=True, still returned so
        callers can surface an upgrade message). Manual links bypass model check.

        Returns a dict with:
          schema_version, cross_product_available, tier_gate, links
        """
        if not is_pro:
            return {
                "schema_version": CROSS_LINK_SCHEMA_VERSION,
                "cross_product_available": True,
                "tier_gate": "pro_required",
                "message": "Cross-product linking requires Pro tier.",
                "links": [],
            }
        with self._db() as conn:
            rows = conn.execute(
                """SELECT target_product, target_id, similarity_score,
                          embedding_model, embedding_dim, link_type, created_at
                   FROM cross_product_links
                   WHERE source_product = ? AND source_id = ?
                   ORDER BY
                     CASE link_type WHEN 'manual' THEN 0 ELSE 1 END,
                     CASE WHEN embedding_model = ? THEN 0 ELSE 1 END,
                     similarity_score DESC
                   LIMIT ?""",
                (source_product, source_id, current_embedding_model, limit),
            ).fetchall()

        links = []
        for row in rows:
            is_stale = (row["link_type"] != "manual" and
                        row["embedding_model"] != current_embedding_model)
            entry: Dict[str, Any] = {
                "target_product": row["target_product"],
                "target_id": row["target_id"],
                "similarity_score": row["similarity_score"],
                "link_type": row["link_type"],
                "created_at": row["created_at"],
                "is_stale": is_stale,
            }
            if is_stale:
                entry["upgrade_message"] = (
                    "This link was created with a different embedding model. "
                    "Re-index both products to refresh cross-product links."
                )
            links.append(entry)

        return {
            "schema_version": CROSS_LINK_SCHEMA_VERSION,
            "cross_product_available": True,
            "tier_gate": "satisfied",
            "links": links,
        }

    def link_session_to_doc(
        self,
        session_id: str,
        doc_id: str,
        vault_id: str,
        link_type: str = "manual",
        similarity_score: Optional[float] = None,
        is_pro: bool = False,
    ) -> Dict[str, Any]:
        """Write a manual cross-product link from a LoreConvo session to a LoreDocs doc.

        Checks vault opt-out. Manual links use embedding_model='manual', dim=0,
        tier_required='free' (manual links are free-tier accessible per design).
        """
        if not is_pro and link_type == "manual":
            # Manual links allowed for all tiers per architecture decision
            pass

        with self._db() as conn:
            vault_row = conn.execute(
                "SELECT cross_link_opt_out FROM vaults WHERE id = ?", (vault_id,)
            ).fetchone()
            if not vault_row:
                return {"ok": False, "reason": "vault not found"}
            if vault_row["cross_link_opt_out"]:
                return {"ok": False, "reason": "vault has cross-linking disabled"}

            doc_row = conn.execute(
                "SELECT 1 FROM documents WHERE id = ? AND deleted = 0", (doc_id,)
            ).fetchone()
            if not doc_row:
                return {"ok": False, "reason": "document not found"}

            self._write_cross_product_link(
                conn,
                source_product="loreconvo",
                source_id=session_id,
                target_product="loredocs",
                target_id=doc_id,
                similarity_score=similarity_score,
                embedding_model="manual",
                embedding_dim=0,
                link_type="manual",
                tier_required="free",
            )

        return {"ok": True, "session_id": session_id, "doc_id": doc_id}

    def cross_link_doc(self, doc_id: str, vault_id: str) -> int:
        """Trigger save-time cross-product linking for a LoreDocs document.

        Queries the LoreConvo sessions.lance index for semantically similar sessions,
        writes up to _CROSS_LINK_CAP links. Pro-only, best-effort (never fails a save).

        Returns count of links written (0 if product unavailable, not Pro, etc).
        Updates documents.last_cross_linked_at.
        """
        import logging as _logging
        log = _logging.getLogger(__name__)

        if get_tier(self.root) != TIER_PRO:
            return 0

        # Check debounce
        now_str = self._now()
        with self._db() as conn:
            row = conn.execute(
                "SELECT last_cross_linked_at FROM documents WHERE id = ?", (doc_id,)
            ).fetchone()
        if row and row["last_cross_linked_at"]:
            try:
                from datetime import datetime as _dt, timezone as _tz
                last = _dt.fromisoformat(row["last_cross_linked_at"].replace("Z", "+00:00"))
                now_dt = _dt.now(_tz.utc)
                if (now_dt - last).total_seconds() < _CROSS_LINK_DEBOUNCE_SECS:
                    return 0
            except Exception:
                pass

        # Discover LoreConvo sessions.lance
        try:
            lc_db = discover_product_db("loreconvo")
        except DiscoveryError:
            return 0
        if lc_db is None:
            return 0

        # Find LoreConvo sessions.lance directory
        lc_lance_dir = lc_db.parent / "sessions.lance"
        if not lc_lance_dir.exists():
            return 0

        # Get the document's text for embedding
        with self._db() as conn:
            vault_row = conn.execute(
                "SELECT cross_link_opt_out FROM vaults WHERE id = ?", (vault_id,)
            ).fetchone()
            if not vault_row or vault_row["cross_link_opt_out"]:
                return 0
            doc_row = conn.execute(
                "SELECT name FROM documents WHERE id = ? AND deleted = 0", (doc_id,)
            ).fetchone()
        if not doc_row:
            return 0

        extracted_path = self.vaults_dir / vault_id / "docs" / doc_id / "extracted.txt"
        doc_text = doc_row["name"]
        if extracted_path.exists():
            try:
                doc_text = extracted_path.read_text(encoding="utf-8", errors="replace")[:2000]
            except Exception:
                pass

        try:
            import lancedb as _lancedb
            from sentence_transformers import SentenceTransformer as _ST

            model = _ST(_CROSS_LINK_EMBEDDING_MODEL)
            q_vec = model.encode(doc_text).tolist()

            lc_lance_db = _lancedb.connect(str(lc_lance_dir))
            table = lc_lance_db.open_table("sessions")
            raw = table.search(
                q_vec, vector_column_name="vector", query_type="vector"
            ).limit(50).to_list()

            # Deduplicate to best distance per session_id
            best: Dict[str, float] = {}
            for r in raw:
                sid = r.get("session_id")
                dist = r.get("_distance", 999.0)
                if not sid:
                    continue
                if dist > _CROSS_LINK_L2_THRESHOLD:
                    continue
                if sid not in best or dist < best[sid]:
                    best[sid] = dist

            # Check LoreConvo DB for session opt-out
            import sqlite3 as _sqlite3
            lc_conn = _sqlite3.connect(f"file:{lc_db}?mode=ro", uri=True)
            lc_conn.row_factory = _sqlite3.Row
            written = 0
            with self._db() as conn:
                for sid, dist in sorted(best.items(), key=lambda x: x[1])[:_CROSS_LINK_CAP]:
                    if written >= _CROSS_LINK_CAP:
                        break
                    row = lc_conn.execute(
                        "SELECT cross_link_opt_out FROM sessions WHERE id = ?", (sid,)
                    ).fetchone()
                    # cross_link_opt_out may not exist on older schemas -- treat as 0
                    if row and row[0]:
                        continue
                    cosine = max(0.0, 1.0 - dist)
                    self._write_cross_product_link(
                        conn,
                        source_product="loredocs",
                        source_id=doc_id,
                        target_product="loreconvo",
                        target_id=sid,
                        similarity_score=round(cosine, 4),
                        embedding_model=_CROSS_LINK_EMBEDDING_MODEL,
                        embedding_dim=_CROSS_LINK_EMBEDDING_DIM,
                        link_type="auto",
                        tier_required="pro",
                    )
                    written += 1
                conn.execute(
                    "UPDATE documents SET last_cross_linked_at = ? WHERE id = ?",
                    (now_str, doc_id),
                )
            lc_conn.close()
            log.debug("cross_link_doc: wrote %d cross-product links for doc %s", written, doc_id)
            return written
        except Exception as exc:
            log.warning("cross_link_doc: unavailable (%s)", type(exc).__name__)
            return 0

    # -------------------------------------------------------------------
    # Suggestions (Phase 2)
    # -------------------------------------------------------------------

    def get_suggestions(self, vault_id: Optional[str] = None,
                        limit: int = 5) -> List[Dict[str, Any]]:
        """Return suggested documents to review or improve.

        Surfaces documents that may need attention:
        - Recently added with no notes (undocumented)
        - No tags (unorganized)
        - Never linked to any other document (isolated)

        Optionally scoped to a single vault.
        """
        suggestions: List[Dict[str, Any]] = []

        with self._db() as conn:
            vault_filter = "AND d.vault_id = ?" if vault_id else ""
            params_base: List[Any] = [vault_id] if vault_id else []

            # 1. Recently added but no notes
            rows = conn.execute(
                f"""SELECT d.id, d.name, d.vault_id, v.name as vault_name,
                           d.category, d.created_at, d.updated_at
                    FROM documents d
                    JOIN vaults v ON d.vault_id = v.id
                    WHERE d.deleted = 0 AND (d.notes IS NULL OR d.notes = '')
                    {vault_filter}
                    ORDER BY d.created_at DESC
                    LIMIT ?""",
                params_base + [limit]
            ).fetchall()
            for row in rows:
                suggestions.append({
                    "reason": "no_notes",
                    "label": "Add notes to describe this document",
                    "doc_id": row["id"],
                    "doc_name": row["name"],
                    "vault_id": row["vault_id"],
                    "vault_name": row["vault_name"],
                    "category": row["category"],
                    "updated_at": row["updated_at"],
                })

            if len(suggestions) >= limit:
                return suggestions[:limit]

            # 2. Documents with no tags
            rows = conn.execute(
                f"""SELECT d.id, d.name, d.vault_id, v.name as vault_name,
                           d.category, d.updated_at
                    FROM documents d
                    JOIN vaults v ON d.vault_id = v.id
                    WHERE d.deleted = 0
                      AND (d.tags IS NULL OR d.tags = '[]')
                    {vault_filter}
                    ORDER BY d.updated_at DESC
                    LIMIT ?""",
                params_base + [limit - len(suggestions)]
            ).fetchall()
            for row in rows:
                if not any(s["doc_id"] == row["id"] for s in suggestions):
                    suggestions.append({
                        "reason": "no_tags",
                        "label": "Tag this document for better discoverability",
                        "doc_id": row["id"],
                        "doc_name": row["name"],
                        "vault_id": row["vault_id"],
                        "vault_name": row["vault_name"],
                        "category": row["category"],
                        "updated_at": row["updated_at"],
                    })

            if len(suggestions) >= limit:
                return suggestions[:limit]

            # 3. Documents never linked to anything
            rows = conn.execute(
                f"""SELECT d.id, d.name, d.vault_id, v.name as vault_name,
                           d.category, d.updated_at
                    FROM documents d
                    JOIN vaults v ON d.vault_id = v.id
                    WHERE d.deleted = 0
                      AND d.id NOT IN (SELECT DISTINCT source_doc_id FROM doc_links)
                    {vault_filter}
                    ORDER BY d.created_at DESC
                    LIMIT ?""",
                params_base + [limit - len(suggestions)]
            ).fetchall()
            for row in rows:
                if not any(s["doc_id"] == row["id"] for s in suggestions):
                    suggestions.append({
                        "reason": "no_links",
                        "label": "Link this document to related documents",
                        "doc_id": row["id"],
                        "doc_name": row["name"],
                        "vault_id": row["vault_id"],
                        "vault_name": row["vault_name"],
                        "category": row["category"],
                        "updated_at": row["updated_at"],
                    })

        return suggestions[:limit]
